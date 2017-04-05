# -*- coding: utf-8 -*-
"""

"""
 
import fnmatch, os,  sys, win32com.client
#import pythoncom
import xlrd
import csv
import subprocess
 
 
wordapp = win32com.client.Dispatch("Word.Application")
 
 
#dirname="C:/BigData/PII_Finder/NameAddressTest/in"
 
mycmd = "C:/Software/xpdfbin-win-3.04/xpdfbin-win-3.04/bin64/pdftotext.exe"
layout = "-layout"
table ="-table"
raw = "-raw"
 
class convert_documents:
   
 
 
    def __init__(self, dirname):
        #    for path, dirs, files in os.walk(sys.argv[1]):
        for path, dirs, files in os.walk(dirname):
        #    print(path,files)
        #    for doc in [os.path.abspath(os.path.join(path, filename)) for filename in files if fnmatch.fnmatch(filename, '*.docx')]:
        #       print("Doc ", doc)
       
            for filename in files:
                if fnmatch.fnmatch(filename, '*.docx'):
                    docx = os.path.abspath(os.path.join(path, filename))
                    print("Found Word Docx ", docx)
                    try:
                        wordapp.Documents.Open(docx)
                        docastxt = docx.rstrip('.docx') + '_DOCX.txt'
#                        wordapp.ActiveDocument.SaveAs(docastxt, FileFormat=win32com.client.constants.wdFormatText)
                        wordapp.ActiveDocument.SaveAs(docastxt, FileFormat=7)
                       
                        wordapp.ActiveDocument.Close()
                    except:
                        print("Bad File:",docx)
                   
                if fnmatch.fnmatch(filename, '*.doc'):
                    doc = os.path.abspath(os.path.join(path, filename))
                    print("Found Word Doc ", doc) 
                    try:
                        wordapp.Documents.Open(doc)
                        docastxt = doc.rstrip('.doc') + '_DOC.txt'
                        wordapp.ActiveDocument.SaveAs(docastxt, FileFormat=7)
                        wordapp.ActiveDocument.Close()
                    except:
                        print("Bad File:",doc)
                        
                 
                    
                if fnmatch.fnmatch(filename, '*.xls'):
                    xlsfile = os.path.abspath(os.path.join(path, filename))
                    print("Found XLS ", xlsfile)  
#                    wrkbk =  xlrd.open_workbook(xlsfile)
                    try:
                        wrkbk =  xlrd.open_workbook(xlsfile)
                        for sheet in wrkbk.sheet_names():
                            print("Sheet ", str(sheet))
                            out_file = xlsfile.rstrip('.xls') + '_SHEET_'+ str(sheet)+'_XLS.txt'
                            sht = wrkbk.sheet_by_name(str(sheet))
                            csvfile = open(out_file, 'wb')
                            writecsv = csv.writer(csvfile, quoting=csv.QUOTE_ALL)
                            #print("Rows:",x1.nrows)
                            for rownum in xrange(sht.nrows):
                                writecsv.writerow(sht.row_values(rownum))
           
                            csvfile.close()
                    except:
                        print("Bad File:",xlsfile)                       
                        
                if fnmatch.fnmatch(filename, '*.xlsx'):
                    xlsxFile = os.path.abspath(os.path.join(path, filename))
                    print("Found XLSX ", xlsxFile)
                     
#                    wrkbk =  xlrd.open_workbook(xlsxFile)
                    try:
                        wrkbk =  xlrd.open_workbook(xlsxFile)
                        for sheet in wrkbk.sheet_names():
                            print("Sheet ", str(sheet))
                            out_file = xlsxFile.rstrip('.xlsx') + '_SHEET_'+ str(sheet)+'_XLSX.txt'
                            sht = wrkbk.sheet_by_name(str(sheet))
                            csvfile = open(out_file, 'wb')
                            writecsv = csv.writer(csvfile, quoting=csv.QUOTE_ALL)
                            #print("Rows:",x1.nrows)
                            for rownum in xrange(sht.nrows):
                                writecsv.writerow(sht.row_values(rownum))
           
                            csvfile.close()
                    except:
                        print("Bad File:",xlsxFile)
                    
                if fnmatch.fnmatch(filename, '*.pdf'):
                    pdfFile = os.path.abspath(os.path.join(path, filename))
                    print("Found PDF ", pdfFile)   
                    out_file = pdfFile.rstrip('.pdf') + '_PDF.txt'
                    p = subprocess.call([mycmd,layout,table,raw,pdfFile,out_file])
                    print(p)
 
                """
               
                from pptx import Presentation
               
                prs = Presentation(path_to_presentation)
               
                # text_runs will be populated with a list of strings,
                # one for each text run in presentation
                text_runs = []
               
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if not shape.has_textframe:
                            continue
                        for paragraph in shape.textframe.paragraphs:
                            for run in paragraph.runs:
                                text_runs.append(run.text)
                """
print(wordapp)
 
       
if __name__ == '__main__':
  args = sys.argv[1:]
  if (len(args) != 1):
    print("Usage: python Convert_WD_XL_PD_To_TXT.py  <Starting Directory for Files> ")
    sys.exit(1)
  dirname = args[0] 
execute = convert_documents(dirname)
 